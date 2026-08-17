#!/usr/bin/env python3
"""
strip-notes.py -- copy a .pptx with the speaker notes removed.

The instructor decks keep their speaker notes; the copies published in this
public repo do not. Run this after editing a deck to refresh the repo copy:

    python 0_System/scripts/strip-notes.py "<source.pptx>" 1_Class/Slides/

Takes one or more source files and an output directory. Each output keeps the
source filename. The source is never modified.

Requires python-pptx:  pip install python-pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
except ImportError:
    sys.exit("python-pptx is not installed.  pip install python-pptx")


def strip(src, out_dir):
    prs = Presentation(src)
    removed = 0

    for slide in prs.slides:
        # has_notes_slide must be checked first -- touching slide.notes_slide
        # CREATES an empty notes slide rather than reporting there isn't one.
        if not slide.has_notes_slide:
            continue
        rids = [rid for rid, rel in slide.part.rels.items()
                if rel.reltype == RT.NOTES_SLIDE]
        for rid in rids:
            slide.part.drop_rel(rid)
            removed += 1

    os.makedirs(out_dir, exist_ok=True)
    dest = os.path.join(out_dir, os.path.basename(src))
    prs.save(dest)

    # Confirm rather than assume: reopen and count.
    check = Presentation(dest)
    left = sum(1 for s in check.slides if s.has_notes_slide)
    print("  %-34s %2d slides, %2d notes removed, %d remaining"
          % (os.path.basename(src), len(check.slides._sldIdLst), removed, left))
    return left


def main():
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    *sources, out_dir = sys.argv[1:]
    print("\nStripping speaker notes -> %s\n" % out_dir)
    remaining = sum(strip(s, out_dir) for s in sources)
    if remaining:
        sys.exit("\nFAILED: %d slide(s) still carry notes." % remaining)
    print("\nDone. Source files unchanged.\n")


if __name__ == "__main__":
    main()
